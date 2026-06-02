#!/usr/bin/env python3
"""Build «Два хакатона: история контрастов» presentation from dva_hakatonа.md plan."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GOLD      = RGBColor(0xF5, 0xA6, 0x23)
ORANGE    = RGBColor(0xE8, 0x6A, 0x17)
RED       = RGBColor(0xC0, 0x39, 0x2B)
GRAY      = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
TEAL      = RGBColor(0x00, 0x96, 0x88)
BLUE_DARK = RGBColor(0x15, 0x3E, 0x63)
LIGHT_GRAY_BG = RGBColor(0xE8, 0xE8, 0xE8)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank


# ── Helpers ────────────────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.15):
    """Add a text box and return (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox, tf


def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=0, space_after=4):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_shape(slide, shape_type, left, top, width, height, fill_color=None,
              line_color=None, line_width=None):
    """Add a shape with optional fill."""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_diagonal_bg(slide, color_left, color_right):
    """Add a diagonal-split background using two triangles."""
    # Left triangle (top-left to bottom-right)
    tri_left = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, SLIDE_W, SLIDE_H)
    tri_left.rotation = 0.0
    tri_left.fill.solid()
    tri_left.fill.fore_color.rgb = color_left
    tri_left.line.fill.background()
    # Position adjustment to cover left half diagonally
    tri_left.left = 0
    tri_left.top = 0

    # Right triangle (bottom-right)
    tri_right = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, SLIDE_W, SLIDE_H)
    tri_right.rotation = 180.0
    tri_right.fill.solid()
    tri_right.fill.fore_color.rgb = color_right
    tri_right.line.fill.background()
    tri_right.left = 0
    tri_right.top = 0


def add_slide_number(slide, number):
    """Add small slide number in bottom-right."""
    add_textbox(slide, 12.5, 7.0, 0.7, 0.4, str(number),
                font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, items, title_color=WHITE,
             bg_color=DARK_BG, text_color=WHITE):
    """Add a card with title and bullet items."""
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill_color=bg_color, line_color=GRAY, line_width=0.5)
    card.shadow.inherit = False
    # Title
    add_textbox(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, title,
                font_size=20, color=title_color, bold=True)
    # Items
    _, tf = add_textbox(slide, left + 0.3, top + 0.8, width - 0.6, height - 1.0,
                        "", font_size=14, color=text_color)
    for item in items:
        add_paragraph(tf, item, font_size=14, color=text_color, space_before=4)


# ── Slide 1: Title ────────────────────────────────────────────────────────
def slide_01():
    slide = prs.slides.add_slide(blank_layout)
    # Diagonal background
    add_diagonal_bg(slide, BLUE_DARK, TEAL)

    # Title
    add_textbox(slide, 1.5, 1.8, 10.3, 1.2,
                "Два хакатона:\nистория контрастов",
                font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                line_spacing=1.25)

    # Subtitle
    add_textbox(slide, 1.5, 3.3, 10.3, 0.8,
                "Московский студенческий DATA-Хакатон 2026  vs  Система Хак: Томск 2026",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Team
    add_textbox(slide, 1.5, 5.0, 10.3, 0.6,
                "Команда: [название / имена участников]",
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


# ── Slide 2: About this talk ──────────────────────────────────────────────
def slide_02():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "О чём этот рассказ",
                font_size=36, color=WHITE, bold=True)

    theses = [
        ("01", "Один хакатон — почти победа\nпри сложном старте", WHITE),
        ("02", "Другой — понятная задача,\nно жёсткие условия", GRAY),
        ("03", "Вывод: что на самом деле\nрешает исход", GOLD),
    ]

    for i, (num, text, clr) in enumerate(theses):
        x = 1.0 + i * 4.0
        # Circle
        circle = add_shape(slide, MSO_SHAPE.OVAL, x + 1.0, 2.2, 1.2, 1.2,
                           fill_color=BLUE_DARK, line_color=TEAL, line_width=2)
        circle_tf = circle.text_frame
        circle_tf.paragraphs[0].text = num
        circle_tf.paragraphs[0].font.size = Pt(24)
        circle_tf.paragraphs[0].font.color.rgb = clr
        circle_tf.paragraphs[0].font.bold = True
        circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle_tf.word_wrap = False

        # Text
        add_textbox(slide, x, 3.7, 3.2, 1.2, text,
                    font_size=16, color=clr, alignment=PP_ALIGN.CENTER)

        # Connecting line (except last)
        if i < 2:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 3.2, 2.75, 0.8, 0.04,
                      fill_color=TEAL)

    add_slide_number(slide, 2)


# ── Slide 3: Why these two ────────────────────────────────────────────────
def slide_03():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Почему именно эти два хакатона?",
                font_size=36, color=WHITE, bold=True)

    # Two cards side by side
    add_card(slide, 1.0, 2.0, 5.0, 4.5, "DATA-Хакатон, Москва",
             ["🧩 Тип задачи: размытая", "⏱️ Время: полное", "👥 Конкуренция: разбросанная"],
             title_color=GOLD, bg_color=RGBColor(0x22, 0x22, 0x3E), text_color=WHITE)

    add_card(slide, 7.3, 2.0, 5.0, 4.5, "Система Хак, Томск",
             ["🎯 Тип задачи: чёткая", "⏱️ Время: –1 день", "👥 Конкуренция: плотная"],
             title_color=TEAL, bg_color=RGBColor(0x22, 0x22, 0x3E), text_color=WHITE)

    # VS
    add_textbox(slide, 5.8, 3.6, 1.7, 0.8, "VS",
                font_size=28, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 3)


# ── Slide 4: First hackathon - mystery task ───────────────────────────────
def slide_04():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8,
                "DATA-Хакатон, Москва — 2026: задача-загадка",
                font_size=32, color=WHITE, bold=True)

    # Task quote - centered, with "spotlight" feel
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, 2.0, 8.3, 1.8,
              fill_color=RGBColor(0x25, 0x25, 0x40),
              line_color=GOLD, line_width=1.5)

    add_textbox(slide, 3.0, 2.2, 7.3, 1.4,
                "«Задача звучала примерно так:\n[расплывчатая формулировка без чётких метрик]»",
                font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Three question marks
    questions = [
        "Что считать результатом?",
        "Как измерять успех?",
        "С чего начать?",
    ]
    for i, q in enumerate(questions):
        y = 4.5 + i * 0.7
        add_textbox(slide, 4.0, y, 5.3, 0.5, f"❓  {q}",
                    font_size=18, color=WHITE)

    add_slide_number(slide, 4)


# ── Slide 5: First hackathon - how we coped ──────────────────────────────
def slide_05():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Как мы выкручивались",
                font_size=36, color=WHITE, bold=True)

    stages = [
        ("1", "Интерпретация", "Сузили задачу сами,\nпереформулировали фокус",
         TEAL, 1.0),
        ("2", "Прототип", "Сделали акцент\nна понятности решения",
         GRAY, 4.8),
        ("3", "Презентация", "Главное оружие —\n3–4 часа на питч",
         GOLD, 8.6),
    ]

    for num, title, desc, clr, x in stages:
        # Number circle
        circle = add_shape(slide, MSO_SHAPE.OVAL, x + 1.1, 2.2, 1.0, 1.0,
                           fill_color=clr, line_color=clr, line_width=2)
        ctf = circle.text_frame
        ctf.paragraphs[0].text = num
        ctf.paragraphs[0].font.size = Pt(28)
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        add_textbox(slide, x, 3.5, 3.2, 0.5, title,
                    font_size=20, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
        # Desc
        add_textbox(slide, x, 4.1, 3.2, 1.2, desc,
                    font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Arrow (except last)
        if num != "3":
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x + 3.2, 2.6, 0.6, 0.3,
                      fill_color=GRAY)

    # Highlight box around stage 3
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.3, 1.8, 3.8, 3.8,
              fill_color=None, line_color=GOLD, line_width=2)

    add_slide_number(slide, 5)


# ── Slide 6: First hackathon - 4th place ──────────────────────────────────
def slide_06():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    # Big number 4
    add_textbox(slide, 5.0, 1.0, 3.3, 2.5, "4",
                font_size=160, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 4.5, 3.5, 4.3, 0.6, "место из 30+ команд",
                font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Quote box
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.0, 4.5, 7.3, 1.5,
              fill_color=RGBColor(0x2A, 0x2A, 0x40),
              line_color=GOLD, line_width=1)

    add_textbox(slide, 3.5, 4.7, 6.3, 1.1,
                "«Подача была лучшей в треке.\nПрезентация вытащила нас туда, куда код\nсам по себе не дотянул бы.»",
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 6)


# ── Slide 7: Second hackathon - clear task ────────────────────────────────
def slide_07():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8,
                "Система Хак, Томск — 2026: ясная задача",
                font_size=32, color=BLACK, bold=True)

    # Task quote - clean, structured
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, 2.0, 8.3, 1.8,
              fill_color=WHITE, line_color=TEAL, line_width=1.5)

    add_textbox(slide, 3.0, 2.2, 7.3, 1.4,
                "«Конкретная задача, измеримые критерии,\nпримеры хорошего результата»",
                font_size=20, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 4.0, 4.5, 5.3, 0.6,
                "✅  Наконец-то понятно, что делать!",
                font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)


# ── Slide 8: Second hackathon - pitfalls ──────────────────────────────────
def slide_08():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Подводные камни",
                font_size=36, color=BLACK, bold=True)

    # Problem 1
    card1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 2.0, 5.3, 4.0,
                      fill_color=WHITE, line_color=RED, line_width=2)
    add_textbox(slide, 1.5, 2.3, 4.3, 0.5, "🚂  –1 день",
                font_size=24, color=RED, bold=True)
    add_textbox(slide, 1.5, 3.2, 4.3, 1.5,
                "Мы приехали на день позже.\nПока другие команды уже\nпилили прототипы,\nмы только заселялись.",
                font_size=16, color=DARK_GRAY)

    # Problem 2
    card2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.0, 2.0, 5.3, 4.0,
                      fill_color=WHITE, line_color=RED, line_width=2)
    add_textbox(slide, 7.5, 2.3, 4.3, 0.5, "🪞  Одинаковые решения",
                font_size=24, color=RED, bold=True)
    add_textbox(slide, 7.5, 3.2, 4.3, 1.5,
                "6 команд из 20 делали\nпочти то же самое.\nВыделиться стало\nневероятно трудно.",
                font_size=16, color=DARK_GRAY)

    add_slide_number(slide, 8)


# ── Slide 9: Second hackathon - timeline ──────────────────────────────────
def slide_09():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8,
                "Что мы успели — и чего не хватило",
                font_size=32, color=BLACK, bold=True)

    days = [
        ("День 1", "Дорога", RED, True),
        ("День 2", "Разработка", TEAL, False),
        ("День 3", "Финал", GOLD, False),
    ]

    for i, (day, event, clr, strikethrough) in enumerate(days):
        x = 1.5 + i * 4.0

        # Day label
        label = f"~~{day}~~" if strikethrough else day
        lbl_color = GRAY if strikethrough else BLACK
        add_textbox(slide, x, 2.0, 3.0, 0.5, label,
                    font_size=22, color=lbl_color, bold=True)

        # Event card
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.8, 3.0, 1.8,
                         fill_color=WHITE if not strikethrough else LIGHT_GRAY_BG,
                         line_color=clr, line_width=2)

        add_textbox(slide, x + 0.3, 3.0, 2.4, 0.5, event,
                    font_size=18, color=clr, bold=True, alignment=PP_ALIGN.CENTER)

        if i == 2:
            add_textbox(slide, x + 0.3, 3.6, 2.4, 0.8,
                        "Первые 3 места:\nкоманды, приехавшие\nзаранее",
                        font_size=12, color=RED, alignment=PP_ALIGN.CENTER)

        # Arrow
        if i < 2:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x + 3.1, 3.4, 0.6, 0.3,
                      fill_color=GRAY)

    # Bottom note
    add_textbox(slide, 2.0, 5.5, 9.3, 0.6,
                "Разница между 3-м и нашим местом — это были детали, на которые не хватило времени",
                font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 9)


# ── Slide 10: Comparison table ────────────────────────────────────────────
def slide_10():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.3, 11.7, 0.7,
                "Сравнение: два хакатона в одной таблице",
                font_size=32, color=WHITE, bold=True)

    # Table
    rows, cols = 7, 3
    table_shape = slide.shapes.add_table(rows, cols,
                                         Inches(0.8), Inches(1.3),
                                         Inches(11.7), Inches(5.5))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.1)
    table.columns[2].width = Inches(4.1)

    data = [
        ["Параметр", "DATA-Хакатон, Москва", "Система Хак, Томск"],
        ["Ясность задачи", "Размытая", "Чёткая"],
        ["Время", "Полное", "–1 день"],
        ["Конкуренция", "Разбросанная\n(разные подходы)", "Плотная\n(одинаковые решения)"],
        ["Наше преимущество", "Интерпретация + подача", "Опыт + уверенность"],
        ["Итог", "4-е место", "Без призов"],
    ]

    # Extra row
    data.append(["Фактор успеха", "Свобода интерпретации", "Время и логистика"])

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_DARK
            elif c == 0:
                p.font.color.rgb = GRAY
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
            else:
                p.font.color.rgb = WHITE
                if r == len(data) - 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40)
                    p.font.bold = True
                    p.font.color.rgb = GOLD
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3E)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x32)

    add_slide_number(slide, 10)


# ── Slide 11: Key takeaways ───────────────────────────────────────────────
def slide_11():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Главные выводы",
                font_size=36, color=WHITE, bold=True)

    takeaways = [
        ("💡", "Размытая задача — это свобода интерпретации. Используйте её.", GOLD),
        ("🎤", "Презентация решает минимум 30% успеха. Не экономьте на ней.", WHITE),
        ("🚩", "Ясная задача — это ловушка одинаковых решений. Ищите отличие.", TEAL),
        ("🗓️", "Время и логистика — такой же ресурс, как навыки. Планируйте их.", ORANGE),
    ]

    for i, (icon, text, clr) in enumerate(takeaways):
        y = 1.7 + i * 1.3
        # Icon
        add_textbox(slide, 1.0, y, 0.7, 0.7, icon,
                    font_size=28, color=clr, alignment=PP_ALIGN.CENTER)
        # Separator line
        add_shape(slide, MSO_SHAPE.RECTANGLE, 1.9, y + 0.15, 0.04, 0.4,
                  fill_color=clr if i != 1 else GRAY)
        # Text
        add_textbox(slide, 2.2, y, 10.0, 0.7, text,
                    font_size=18, color=clr)

    add_slide_number(slide, 11)


# ── Slide 12: Final ───────────────────────────────────────────────────────
def slide_12():
    slide = prs.slides.add_slide(blank_layout)
    # Diagonal background (swapped colors)
    add_diagonal_bg(slide, TEAL, BLUE_DARK)

    # Quote
    add_textbox(slide, 1.5, 1.0, 10.3, 1.5,
                "«Хакатон — это не только про код»",
                font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.5, 2.7, 10.3, 1.2,
                "Это про людей, время, подачу и умение видеть\nвозможности там, где другие видят хаос",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.0, 4.3, 3.3, 0.03,
              fill_color=GOLD)

    # Summary
    add_textbox(slide, 1.5, 4.7, 10.3, 0.7,
                "Первый хакатон вытащили подачей и интерпретацией.\nВторой — проиграли времени и конкуренции.\nНо в обоих случаях мы выросли.",
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Contacts
    add_textbox(slide, 1.5, 6.2, 10.3, 0.5,
                "Контакты: Telegram @...  |  GitHub @...",
                font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 12)


# ── Build ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    slide_01()
    slide_02()
    slide_03()
    slide_04()
    slide_05()
    slide_06()
    slide_07()
    slide_08()
    slide_09()
    slide_10()
    slide_11()
    slide_12()

    output_path = "/Users/elvsevolod/Desktop/Два_хакатона_история_контрастов.pptx"
    prs.save(output_path)
    print(f"✅ Saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
